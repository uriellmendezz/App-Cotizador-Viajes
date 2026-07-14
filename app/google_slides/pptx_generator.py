import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return RGBColor(*(int(hex_str[i:i+2], 16) for i in (0, 2, 4)))

def add_header(slide, title_text, logo_path, accent_hex, text_color_hex="#ffffff"):
    """Adds a standard header banner on slides 2, 3, 4, 5."""
    # Top banner background
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = hex_to_rgb(accent_hex)
    banner.line.color.rgb = hex_to_rgb(accent_hex)
    
    # Title Text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Montserrat'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(text_color_hex)
    p.alignment = PP_ALIGN.LEFT
    
    # Logo
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, Inches(10.8), Inches(0.25), width=Inches(2.0), height=Inches(0.7))
        except Exception:
            pass

def format_bold_label_pptx(text_frame, label, value, font_size=12, font_color_rgb=RGBColor(51, 51, 51)):
    """Appends a paragraph with bold label and normal value."""
    p = text_frame.add_paragraph()
    p.space_after = Pt(4)
    run_label = p.add_run()
    run_label.text = label + ": "
    run_label.font.name = 'Montserrat'
    run_label.font.bold = True
    run_label.font.size = Pt(font_size)
    run_label.font.color.rgb = font_color_rgb
    
    run_val = p.add_run()
    run_val.text = str(value)
    run_val.font.name = 'Montserrat'
    run_val.font.bold = False
    run_val.font.size = Pt(font_size)
    run_val.font.color.rgb = font_color_rgb

def set_cell_font(cell, text, bold=False, size=11, color_rgb=RGBColor(51, 51, 51)):
    cell.text_frame.text = text
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.name = 'Montserrat'
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color_rgb

def generate_pptx(quote, output_path):
    """
    Generates a premium native editable PPTX presentation using python-pptx.
    Matches the exact 6-step ordering specified in to-do.md.
    """
    prs = Presentation()
    # 16:9 widescreen dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Fetch brand customization settings
    hoteles_raw = quote.get('hoteles', [])
    moneda = quote.get('moneda')
    if not moneda:
        for h in hoteles_raw:
            if h.get("nombre") in ("METADATA_COTIZACION", "METADATA_PRESUPUESTO_RAPIDO"):
                moneda = h.get("moneda")
                break
    if not moneda:
        moneda = "USD"
    quote["moneda"] = moneda
    quote["hoteles"] = [h for h in hoteles_raw if h.get("nombre") not in ("METADATA_COTIZACION", "METADATA_PRESUPUESTO_RAPIDO")]

    colores = quote.get('colores') or quote.get('agencia_colores') or ["#ff545d", "#343434", "#f79646"]
    accent_hex = colores[0]
    secondary_hex = colores[1] if len(colores) > 1 else "#343434"
    logo_path = quote.get('agencia_logo')
    agency_name = quote.get('agencia_nombre', "One Trip Giordano")
    
    # ----------------------------------------------------
    # SLIDE 1: Cover Slide (Portada)
    # ----------------------------------------------------
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    
    # Background Accent block (left stripe)
    bg_shape = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(5.0), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = hex_to_rgb(accent_hex)
    bg_shape.line.fill.background()
    
    # Logo on cover
    if logo_path and os.path.exists(logo_path):
        try:
            slide1.shapes.add_picture(logo_path, Inches(0.5), Inches(0.5), width=Inches(3.5), height=Inches(1.2))
        except Exception:
            pass
            
    # Brand Name on Cover
    brandBox = slide1.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(4), Inches(0.8))
    p_brand = brandBox.text_frame.paragraphs[0]
    p_brand.text = agency_name.upper()
    p_brand.font.name = 'Montserrat'
    p_brand.font.size = Pt(16)
    p_brand.font.bold = True
    p_brand.font.color.rgb = RGBColor(255, 255, 255)
    
    # Destination & Propuesta Title (large, centered right side)
    titleBox = slide1.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(7.3), Inches(3.0))
    tf1 = titleBox.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "PROPUESTA DE VIAJE"
    p1.font.name = 'Montserrat'
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = hex_to_rgb(secondary_hex)
    
    p2 = tf1.add_paragraph()
    p2.text = quote.get('destino', "DESTINO").upper()
    # Montserrat Black equivalent: bold, large tracking
    p2.font.name = 'Montserrat'
    p2.font.size = Pt(54)
    p2.font.bold = True
    p2.font.color.rgb = hex_to_rgb(accent_hex)
    p2.space_before = Pt(10)
    
    # Passenger and details
    detailsBox = slide1.shapes.add_textbox(Inches(5.5), Inches(4.5), Inches(7.3), Inches(2.0))
    tf_details = detailsBox.text_frame
    tf_details.word_wrap = True
    
    p_pass = tf_details.paragraphs[0]
    p_pass.text = f"Presentado a: {quote.get('nombre_pax', 'Pasajero')}"
    p_pass.font.name = 'Montserrat'
    p_pass.font.size = Pt(18)
    p_pass.font.bold = True
    p_pass.font.color.rgb = RGBColor(51, 51, 51)
    
    p_pax_count = tf_details.add_paragraph()
    p_pax_count.text = f"Cantidad de Pasajeros: {quote.get('cantidad_pasajeros', 1)} | Origen: {quote.get('origen', 'Córdoba')}"
    p_pax_count.font.name = 'Montserrat'
    p_pax_count.font.size = Pt(14)
    p_pax_count.font.color.rgb = RGBColor(102, 102, 102)
    p_pax_count.space_before = Pt(10)
    
    p_date = tf_details.add_paragraph()
    p_date.text = f"Fecha de Salida: {quote.get('fecha_salida', '')}"
    p_date.font.name = 'Montserrat'
    p_date.font.size = Pt(14)
    p_date.font.color.rgb = RGBColor(102, 102, 102)

    # ----------------------------------------------------
    # SLIDE 2: Servicios Cotizados Table (Minimalist style)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    add_header(slide2, "SERVICIOS COTIZADOS EN LA PROPUESTA", logo_path, accent_hex)
    
    # Table details
    rows = 5
    cols = 2
    left = Inches(1.0)
    top = Inches(2.0)
    width = Inches(11.3)
    height = Inches(4.2)
    
    table_shape = slide2.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(8.3)
    
    # Header cells
    set_cell_font(table.cell(0, 0), "TIPO DE PRESTACIÓN", bold=True, size=13, color_rgb=hex_to_rgb(accent_hex))
    set_cell_font(table.cell(0, 1), "DETALLE Y COBERTURA (SIN DESGLOSE ECONÓMICO)", bold=True, size=13, color_rgb=hex_to_rgb(accent_hex))
    
    # Table data
    destino = quote.get('destino', '')
    cant_pax = quote.get('cantidad_pasajeros', 1)
    origen = quote.get('origen', 'Córdoba')
    hoteles = quote.get('hoteles', [])
    hotel_names = ", ".join([h.get('nombre') for h in hoteles if h.get('nombre')])
    noches = quote.get('noches_alojamiento', '7 noches')
    
    set_cell_font(table.cell(1, 0), "Aéreo", bold=True, size=12)
    pax_str = "un pasajero adulto" if cant_pax == 1 else f"{cant_pax} pasajeros adultos"
    set_cell_font(table.cell(1, 1), f"Vuelos desde {origen} a {destino} para {pax_str}.", size=12)
    
    set_cell_font(table.cell(2, 0), "Terrestre (Hotel)", bold=True, size=12)
    set_cell_font(table.cell(2, 1), f"Estadía en hotel ({hotel_names}) por {noches}.", size=12)
    
    set_cell_font(table.cell(3, 0), "Logística", bold=True, size=12)
    set_cell_font(table.cell(3, 1), "Traslados de llegada (aeropuerto-hotel) y salida (hotel-aeropuerto) en destino.", size=12)
    
    set_cell_font(table.cell(4, 0), "Asistencia", bold=True, size=12)
    set_cell_font(table.cell(4, 1), "Seguro de asistencia médica al viajero con cobertura integral.", size=12)

    # ----------------------------------------------------
    # SLIDE 3: Aéreos (Outbound and Inbound flights)
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    add_header(slide3, "DETALLE DE TRAMOS AÉREOS", logo_path, accent_hex)
    
    # Outbound flight card (left column)
    outbound_card = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.8), Inches(5.7), Inches(5.0)
    )
    outbound_card.fill.solid()
    outbound_card.fill.fore_color.rgb = RGBColor(245, 247, 250)
    outbound_card.line.color.rgb = RGBColor(220, 225, 230)
    
    outbound_title = slide3.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(5.3), Inches(0.5))
    p_o = outbound_title.text_frame.paragraphs[0]
    p_o.text = "VUELO DE IDA ✈"
    p_o.font.name = 'Montserrat'
    p_o.font.size = Pt(18)
    p_o.font.bold = True
    p_o.font.color.rgb = hex_to_rgb(accent_hex)
    
    outbound_info = slide3.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(5.3), Inches(1.5))
    tf_out = outbound_info.text_frame
    tf_out.word_wrap = True
    tf_out.margin_left = tf_out.margin_right = tf_out.margin_top = tf_out.margin_bottom = 0
    format_bold_label_pptx(tf_out, "Fecha de Vuelo", quote.get('fecha_vuelo_ida', quote.get('fecha_salida', '')))
    
    img_ida = quote.get('img_vuelo_ida')
    if img_ida and os.path.exists(img_ida):
        try:
            slide3.shapes.add_picture(img_ida, Inches(0.9), Inches(3.8), width=Inches(5.3), height=Inches(2.6))
        except Exception:
            pass
            
    # Inbound flight card (right column)
    inbound_card = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.7), Inches(5.0)
    )
    inbound_card.fill.solid()
    inbound_card.fill.fore_color.rgb = RGBColor(245, 247, 250)
    inbound_card.line.color.rgb = RGBColor(220, 225, 230)
    
    inbound_title = slide3.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.5))
    p_i = inbound_title.text_frame.paragraphs[0]
    p_i.text = "VUELO DE RETORNO ✈"
    p_i.font.name = 'Montserrat'
    p_i.font.size = Pt(18)
    p_i.font.bold = True
    p_i.font.color.rgb = hex_to_rgb(accent_hex)
    
    inbound_info = slide3.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.3), Inches(1.5))
    tf_in = inbound_info.text_frame
    tf_in.word_wrap = True
    tf_in.margin_left = tf_in.margin_right = tf_in.margin_top = tf_in.margin_bottom = 0
    format_bold_label_pptx(tf_in, "Fecha de Vuelo", quote.get('fecha_vuelo_vuelta', ''))
    
    img_vuelta = quote.get('img_vuelo_vuelta')
    if img_vuelta and os.path.exists(img_vuelta):
        try:
            slide3.shapes.add_picture(img_vuelta, Inches(7.1), Inches(3.8), width=Inches(5.3), height=Inches(2.6))
        except Exception:
            pass

    # ----------------------------------------------------
    # SLIDE 4: Terrestre (Hotels with descriptions and specific prices)
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    add_header(slide4, "OPCIONES DE ALOJAMIENTO SUGERIDAS", logo_path, accent_hex)
    
    num_hoteles = len(hoteles)
    base_label = quote.get('base_habitacion', 'Doble')
    
    if num_hoteles == 1:
        # Full slide layout for 1 hotel
        hotel = hoteles[0]
        card = slide4.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.8), Inches(11.93), Inches(5.0)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(245, 247, 250)
        card.line.color.rgb = RGBColor(220, 225, 230)
        
        # Details left
        detailBox = slide4.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(6.0), Inches(4.6))
        tf_h = detailBox.text_frame
        tf_h.word_wrap = True
        tf_h.margin_left = tf_h.margin_right = tf_h.margin_top = tf_h.margin_bottom = 0
        
        p_name = tf_h.paragraphs[0]
        p_name.text = hotel.get('nombre', 'Hotel').upper()
        p_name.font.name = 'Montserrat'
        p_name.font.size = Pt(24)
        p_name.font.bold = True
        p_name.font.color.rgb = hex_to_rgb(accent_hex)
        p_name.space_after = Pt(6)
        
        format_bold_label_pptx(tf_h, "Categoría", f"{hotel.get('estrellas', '')}")
        format_bold_label_pptx(tf_h, "Régimen", hotel.get('regimen', ''))
        format_bold_label_pptx(tf_h, "Tipo de Habitación", hotel.get('habitacion', ''))
        format_bold_label_pptx(tf_h, "Descripción", hotel.get('descripcion', ''))
        
        # Price in card
        p_pr = tf_h.add_paragraph()
        p_pr.space_before = Pt(12)
        run_pr_label = p_pr.add_run()
        run_pr_label.text = "Tarifa Final Paquete: "
        run_pr_label.font.name = 'Montserrat'
        run_pr_label.font.bold = True
        run_pr_label.font.size = Pt(14)
        
        run_pr_val = p_pr.add_run()
        run_pr_val.text = f"{moneda} {hotel.get('costo', 0.0):,.2f}"
        run_pr_val.font.name = 'Montserrat'
        run_pr_val.font.bold = True
        run_pr_val.font.size = Pt(20)
        run_pr_val.font.color.rgb = hex_to_rgb(secondary_hex)
        
        p_pax = tf_h.add_paragraph()
        run_pax = p_pax.add_run()
        # Calculate price per person
        tot = hotel.get('costo', 0.0)
        per = tot / cant_pax if cant_pax > 0 else tot
        run_pax.text = f"{moneda} {per:,.2f} por persona en Base {base_label}"
        run_pax.font.name = 'Montserrat'
        run_pax.font.size = Pt(11)
        run_pax.font.color.rgb = RGBColor(102, 102, 102)
        
        # Image right
        img_hotel = hotel.get('imagen1') or hotel.get('imagen') or hotel.get('imagen2') or hotel.get('imagen3')
        if img_hotel and os.path.exists(img_hotel):
            try:
                slide4.shapes.add_picture(img_hotel, Inches(7.3), Inches(2.1), width=Inches(5.0), height=Inches(4.4))
            except Exception:
                pass
                
    elif num_hoteles > 1:
        # Multiple hotels grid (2 to 4 hotels)
        max_col_width = 11.93
        gap = 0.3
        width_per_hotel = (max_col_width - (gap * (num_hoteles - 1))) / num_hoteles
        
        for idx, hotel in enumerate(hoteles):
            left_pos = 0.7 + idx * (width_per_hotel + gap)
            card = slide4.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(1.8), Inches(width_per_hotel), Inches(5.0)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(245, 247, 250)
            card.line.color.rgb = RGBColor(220, 225, 230)
            
            # Text container
            infoBox = slide4.shapes.add_textbox(Inches(left_pos + 0.15), Inches(2.0), Inches(width_per_hotel - 0.3), Inches(2.8))
            tf_h = infoBox.text_frame
            tf_h.word_wrap = True
            tf_h.margin_left = tf_h.margin_right = tf_h.margin_top = tf_h.margin_bottom = 0
            
            p_name = tf_h.paragraphs[0]
            p_name.text = hotel.get('nombre', 'Hotel')
            p_name.font.name = 'Montserrat'
            p_name.font.size = Pt(13)
            p_name.font.bold = True
            p_name.font.color.rgb = hex_to_rgb(accent_hex)
            p_name.space_after = Pt(4)
            
            format_bold_label_pptx(tf_h, "Cat", hotel.get('estrellas', ''), font_size=9)
            format_bold_label_pptx(tf_h, "Régimen", hotel.get('regimen', ''), font_size=9)
            format_bold_label_pptx(tf_h, "Hab", hotel.get('habitacion', ''), font_size=9)
            
            # Description
            desc = hotel.get('descripcion', '')
            if len(desc) > 50:
                desc = desc[:47] + "..."
            format_bold_label_pptx(tf_h, "Desc", desc, font_size=8)
            
            # Pricing in card
            tot = hotel.get('costo', 0.0)
            per = tot / cant_pax if cant_pax > 0 else tot
            
            p_pr = tf_h.add_paragraph()
            p_pr.space_before = Pt(6)
            run_pr = p_pr.add_run()
            run_pr.text = f"Total: {moneda} {tot:,.2f}\n"
            run_pr.font.name = 'Montserrat'
            run_pr.font.bold = True
            run_pr.font.size = Pt(11)
            run_pr.font.color.rgb = hex_to_rgb(secondary_hex)
            
            run_pax = p_pr.add_run()
            run_pax.text = f"{moneda} {per:,.2f} / pax ({base_label})"
            run_pax.font.name = 'Montserrat'
            run_pax.font.size = Pt(9)
            run_pax.font.color.rgb = RGBColor(102, 102, 102)
            
            # Image bottom
            img_hotel = hotel.get('imagen1') or hotel.get('imagen') or hotel.get('imagen2') or hotel.get('imagen3')
            if img_hotel and os.path.exists(img_hotel):
                try:
                    slide4.shapes.add_picture(
                        img_hotel, 
                        Inches(left_pos + 0.15), 
                        Inches(4.8), 
                        width=Inches(width_per_hotel - 0.3), 
                        height=Inches(1.8)
                    )
                except Exception:
                    pass

    # ----------------------------------------------------
    # SLIDE 5: Cierre Comercial & Legales
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    add_header(slide5, "RESUMEN DE PROPUESTA Y LEGALES", logo_path, accent_hex)
    
    # Highlighted Pricing Summary
    price_card = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.8), Inches(11.93), Inches(2.2)
    )
    price_card.fill.solid()
    price_card.fill.fore_color.rgb = hex_to_rgb(accent_hex)
    price_card.line.fill.background()
    
    pBox = slide5.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.33), Inches(1.6))
    tf_p = pBox.text_frame
    tf_p.word_wrap = True
    tf_p.margin_left = tf_p.margin_right = tf_p.margin_top = tf_p.margin_bottom = 0
    
    costo_total = quote.get('costo_total', 0.0)
    precio_persona = quote.get('precio_persona', 0.0)
    
    p_lbl = tf_p.paragraphs[0]
    p_lbl.text = f"PROPUESTA COMERCIAL - {quote.get('destino','').upper()} (DESDE)"
    p_lbl.font.name = 'Montserrat'
    p_lbl.font.size = Pt(14)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = RGBColor(230, 230, 230)
    p_lbl.alignment = PP_ALIGN.CENTER
    
    p_tot = tf_p.add_paragraph()
    p_tot.alignment = PP_ALIGN.CENTER
    p_tot.space_before = Pt(6)
    r_tot = p_tot.add_run()
    r_tot.text = f"Total {moneda} {costo_total:,.2f}  |  "
    r_tot.font.name = 'Montserrat'
    r_tot.font.size = Pt(28)
    r_tot.font.bold = True
    r_tot.font.color.rgb = RGBColor(255, 255, 255)
    
    r_ind = p_tot.add_run()
    r_ind.text = f"Por Persona {moneda} {precio_persona:,.2f} en Base {base_label}"
    r_ind.font.name = 'Montserrat'
    r_ind.font.size = Pt(18)
    r_ind.font.bold = True
    r_ind.font.color.rgb = RGBColor(240, 240, 240)
    
    # Legal footnote at bottom
    legalBox = slide5.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(11.93), Inches(2.5))
    tf_l = legalBox.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = 0
    
    p_leg = tf_l.paragraphs[0]
    p_leg.alignment = PP_ALIGN.JUSTIFY
    
    legal_text = (
        f"*TASAS TURÍSTICAS OBLIGATORIAS EN AQUELLOS DESTINOS QUE ASÍ LO REQUIERAN. "
        f"*ES DEBER DEL PASAJERO CONOCER LA DOCUMENTACIÓN REQUERIDA. CONSULTÁ LOS REQUISITOS DE INGRESO A CADA DESTINO. "
        f"*LAS TARIFAS COTIZADAS NO EMITIDAS O RESERVADAS PUEDEN SUFRIR CAMBIOS SIN PREVIO AVISO. "
        f"*{agency_name.upper()}, EN SU CARÁCTER DE INTERMEDIARIO Y DISPUESTO A REALIZAR LO QUE ESTÉ A SU ALCANCE, "
        f"NO SE HACE RESPONSABLE POR POLÍTICAS DE REPROGRAMACIONES O CANCELACIONES DE CADA LÍNEA AÉREA U OTROS PROVEEDORES "
        f"DE SERVICIOS TURÍSTICOS ANTE CASOS EVENTUALES."
    )
    
    run_leg = p_leg.add_run()
    run_leg.text = legal_text
    run_leg.font.name = 'Montserrat'
    run_leg.font.size = Pt(8.5)
    run_leg.font.italic = True
    run_leg.font.color.rgb = RGBColor(128, 128, 128)
    
    # Save presentation
    prs.save(output_path)
